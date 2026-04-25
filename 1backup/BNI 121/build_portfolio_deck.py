"""HopeTech Portfolio Deck.

Sources used (fetched 25 Apr 2026):
  - bettroi.com           — Bettroi FZE product portfolio + Dubai HQ
  - drmhope.com           — DrM Hope Softwares: 100+ AI projects, 98% client sat
  - hopephoenix.in        — HopePhoenix India×UAE JV: 30+ products, 100+ team
  - hopetech.me           — Founders + 4-company umbrella structure
  - github.com/chatgptnotes (133 repos) — engineering proof points
  - BNI CRM Supabase database — client list

Run:  python3 build_portfolio_deck.py
Out:  HopeTech_Portfolio_2026-04-25.pptx (next to this file)
"""

from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Brand palette ────────────────────────────────────────────────────────────
NAVY     = RGBColor(0x0F, 0x17, 0x2A)
INDIGO   = RGBColor(0x4F, 0x46, 0xE5)
BLUE     = RGBColor(0x25, 0x63, 0xEB)
VIOLET   = RGBColor(0x7C, 0x3A, 0xED)
PINK     = RGBColor(0xEC, 0x48, 0x99)
GREEN    = RGBColor(0x10, 0xB9, 0x81)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
SLATE    = RGBColor(0x47, 0x55, 0x69)
SLATE_LT = RGBColor(0x94, 0xA3, 0xB8)
BG       = RGBColor(0xF8, 0xFA, 0xFC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide canvas (16:9, 13.33 × 7.5 in) ──────────────────────────────────────
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ── Helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line if line else fill
    s.line.width = Pt(0.25 if line else 0)
    s.shadow.inherit = False
    return s

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Inter',
             tracking=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.name = font; f.size = Pt(size); f.bold = bold; f.color.rgb = color
    if tracking is not None:
        # tracking via XML — letterSpacing in 1/100 pt; rough proxy here
        from pptx.oxml.ns import qn
        rPr = r._r.get_or_add_rPr()
        rPr.set('spc', str(int(tracking * 100)))
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=13, color=NAVY, gap=4):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = Pt(0)
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        # leading bullet glyph
        rb = p.add_run(); rb.text = '▸  '
        rb.font.name = 'Inter'; rb.font.size = Pt(size); rb.font.bold = True
        rb.font.color.rgb = BLUE
        rt = p.add_run(); rt.text = line
        rt.font.name = 'Inter'; rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return tb

def page_chrome(slide, title, subtitle=None, page_no=None):
    """Render the white background, top accent bar, title, and footer."""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, BG)
    # Top accent gradient strip (faked by 3 stacked rects)
    add_rect(slide, 0, 0, Inches(4.5), Inches(0.18), BLUE)
    add_rect(slide, Inches(4.5), 0, Inches(4.5), Inches(0.18), VIOLET)
    add_rect(slide, Inches(9), 0, Inches(4.4), Inches(0.18), PINK)
    add_text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
             title, size=30, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=13, color=SLATE)
    # Footer
    add_text(slide, Inches(0.6), Inches(7.05), Inches(8), Inches(0.3),
             'HopeTech Portfolio · ' + date.today().strftime('%d %b %Y'),
             size=9, color=SLATE_LT)
    if page_no is not None:
        add_text(slide, Inches(12.2), Inches(7.05), Inches(0.6), Inches(0.3),
                 f'{page_no:02d}', size=9, color=SLATE_LT, align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
# Decorative gradient blobs (rounded rectangles)
blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-2), Inches(7), Inches(7))
blob.fill.solid(); blob.fill.fore_color.rgb = VIOLET
blob.line.fill.background(); blob.shadow.inherit = False
blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(4), Inches(6), Inches(6))
blob2.fill.solid(); blob2.fill.fore_color.rgb = INDIGO
blob2.line.fill.background(); blob2.shadow.inherit = False

add_text(s, Inches(0.8), Inches(1.2), Inches(11), Inches(0.4),
         'HOPETECH GROUP · 2026', size=11, bold=True, color=SLATE_LT, tracking=3)
add_text(s, Inches(0.8), Inches(1.85), Inches(12), Inches(2.4),
         'AI · Healthcare · Industrial Automation',
         size=52, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.85), Inches(11), Inches(1.2),
         'A portfolio of products, projects and clients across India and the UAE — '
         'built by Dr. BK Murali and Biji Thomas, delivered through DrM Hope Softwares, '
         'Bettroi FZE, HopePhoenix and Betser Life.',
         size=18, color=RGBColor(0xCB, 0xD5, 0xE1))

# Founder badges
def founder_badge(x, name, role, place):
    add_rect(s, x, Inches(5.5), Inches(5.6), Inches(1.4),
             RGBColor(0x1E, 0x29, 0x3B), line=RGBColor(0x33, 0x41, 0x55))
    add_text(s, x + Inches(0.3), Inches(5.65), Inches(5), Inches(0.4),
             name, size=18, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), Inches(6.05), Inches(5), Inches(0.35),
             role, size=11, color=RGBColor(0xA5, 0xB4, 0xFC))
    add_text(s, x + Inches(0.3), Inches(6.4), Inches(5), Inches(0.4),
             place, size=11, color=SLATE_LT)

founder_badge(Inches(0.8),  'Dr. BK Murali',  'Founder, DrM Hope · CTO, Bettroi · Director, Betser Life',  'Nagpur, India · BNI Diorite')
founder_badge(Inches(6.95), 'Biji Thomas',    'CEO, Bettroi FZE · Co-Founder, HopePhoenix · Betser Life',  'Dubai, UAE')


# ═════════════════════════════════════════════════════════════════════════════
# Slide 2 — Four-company structure
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'The Group', 'Four operating companies, one shared mission.', 2)

cards = [
    ('Bettroi FZE',      'Dubai Silicon Oasis · UAE', 'AI-driven business process automation & consulting',
     'Better business through AI automation', BLUE),
    ('DrM Hope Softwares','Nagpur · India',           'Healthcare AI, generative AI, fintech & autonomous agents',
     'Founded 2010 · 100+ AI projects · 98% satisfaction', INDIGO),
    ('HopePhoenix',      'India × UAE Joint Venture',  'Where Medicine Meets Machines',
     '30+ products · 100+ team · India & Middle East', VIOLET),
    ('Betser Life',      'Thiruvananthapuram · India', 'Preventive health platform · digital + applied AI',
     '750+ healthcare users across 12 countries', PINK),
]
for i, (name, loc, line, foot, color) in enumerate(cards):
    col, row = i % 2, i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.65 + row * 2.65)
    add_rect(s, x, y, Inches(6), Inches(2.4), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(0.18), Inches(2.4), color)
    add_text(s, x + Inches(0.45), y + Inches(0.25), Inches(5.5), Inches(0.45),
             name, size=22, bold=True, color=NAVY)
    add_text(s, x + Inches(0.45), y + Inches(0.75), Inches(5.5), Inches(0.3),
             loc.upper(), size=9, bold=True, color=color, tracking=2)
    add_text(s, x + Inches(0.45), y + Inches(1.15), Inches(5.5), Inches(0.7),
             line, size=14, color=SLATE)
    add_text(s, x + Inches(0.45), y + Inches(1.85), Inches(5.5), Inches(0.4),
             foot, size=10, color=SLATE_LT)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 3 — Stats / scoreboard
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Track Record', 'Numbers compiled from public site disclosures and our GitHub footprint.', 3)

stats = [
    ('100+',   'AI projects delivered',        'DrM Hope Softwares', BLUE),
    ('98%',    'Client satisfaction',          'self-reported',       GREEN),
    ('133',    'Public GitHub repositories',   'github.com/chatgptnotes', INDIGO),
    ('30+',    'Live products in market',      'HopePhoenix portfolio', VIOLET),
    ('750+',   'Healthcare users',             'Betser Life · 12 countries', PINK),
    ('100+',   'Engineering team',             'across India & UAE',  AMBER),
]
for i, (big, label, src, color) in enumerate(stats):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.7 + row * 2.55)
    add_rect(s, x, y, Inches(3.95), Inches(2.3), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(3.95), Inches(0.12), color)
    add_text(s, x + Inches(0.3), y + Inches(0.35), Inches(3.5), Inches(0.95),
             big, size=46, bold=True, color=color)
    add_text(s, x + Inches(0.3), y + Inches(1.4), Inches(3.5), Inches(0.4),
             label, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.3), y + Inches(1.8), Inches(3.5), Inches(0.4),
             src, size=10, color=SLATE_LT)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 4 — Healthcare products
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Healthcare AI · Hospital Software',
            'Software running real hospitals + AI tools for clinicians.', 4)

healthcare = [
    ('AdamRIT',                'Modern HMIS · ESIC patient mapping system'),
    ('OnePatientOneDoctor',    'AI-led patient–doctor continuity ecosystem'),
    ('Aide',                   'Hospital neural network — clinical decision support'),
    ('IHMS / HMS+',            'Hospital management + legacy-system integration'),
    ('Doctors Digital Office', 'OPD / IPD / ICU module suite'),
    ('Yellow Fever vacc. ledger', 'Government-grade vaccine certificate platform'),
    ('Clinivoice',             'Voice-driven clinic notes'),
    ('NABH Online',            'NABH compliance & audit automation'),
    ('Raftaar Health',         'Patient transport & ambulance ops'),
    ('Justdial Lead Mgmt',     'Justdial → Hope Hospital lead pipeline'),
    ('DigiHealthTwin',         'Diabetes digital-twin app'),
    ('OneScanOneLife',         'Mass screening / camp orchestration'),
]
for i, (name, sub) in enumerate(healthcare):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.55 + row * 1.32)
    add_rect(s, x, y, Inches(3.95), Inches(1.18), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y + Inches(0.18), Inches(0.05), Inches(0.82), GREEN)
    add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(3.65), Inches(0.42),
             name, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(3.65), Inches(0.5),
             sub, size=10, color=SLATE)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 5 — Industrial Automation products
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Industrial Automation · SCADA · Energy',
            'Software for the shop floor, the substation and the pipeline.', 5)

industrial = [
    ('FactoryPulse',         'Manufacturing execution system (MES)'),
    ('Ampris',               'Electrical SCADA & grid orchestration'),
    ('FlowNexus',            'Oil-and-gas pipeline telemetry'),
    ('NexaProc',             'Process automation & batch control'),
    ('PLC Pilot',            'AI-generated PLC code (Siemens / Rockwell / Schneider)'),
    ('AutoPanel Design',     'Voice-driven control-panel CAD'),
    ('EcoLogic PLC Asst.',   'Tag mapping + alarm engineering copilot'),
    ('SCADA.work',           'Browser-first SCADA HMI platform'),
    ('Manikaran Dam Safety', 'IoT compliance & reporting for Indian dams'),
    ('Sterling Tender Gen.', 'Auto tender doc generator (Sterling Electricals)'),
    ('IOC',                  'Indian Oil Corp pilot work'),
    ('Vision Claw',          'Vision AI for industrial pick-and-place'),
]
for i, (name, sub) in enumerate(industrial):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.55 + row * 1.32)
    add_rect(s, x, y, Inches(3.95), Inches(1.18), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y + Inches(0.18), Inches(0.05), Inches(0.82), AMBER)
    add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(3.65), Inches(0.42),
             name, size=14, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.6), Inches(3.65), Inches(0.5),
             sub, size=10, color=SLATE)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 6 — Enterprise AI products
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Enterprise AI · Agents · Tools',
            'Bettroi platform stack and standalone AI products.', 6)

enterprise = [
    ('AgentX',          'Autonomous AI agents per business function'),
    ('APX',             'AI Command Center for the enterprise'),
    ('EngageX',         'Omnichannel customer-engagement hub'),
    ('ConnectX',        'Omnichannel communication center'),
    ('Privata',         'Offline / edge AI assistant for sensitive data'),
    ('AISHU',           'Education / tutoring OS'),
    ('AgentSDR',        'Outbound sales-rep AI agent'),
    ('ZeroRiskAgent',   'Sanctions-screening & risk AI'),
    ('Linkist NFC',     'Smart NFC business cards with AI back-office'),
    ('IndiaMonitor',    'Real-time India intelligence dashboard'),
    ('MovingEstimator', 'Video-scan AI for moving / packing estimates'),
    ('PrewedAI',        'Pre-wedding AI photography'),
    ('Hisab',           'Proposal & project tracker (hisab.work)'),
    ('PulseOfProject',  'Project heartbeat dashboard (pulseofproject.com)'),
    ('Ironbark',        'Self-improving learning loop for Claude Code (open-source)'),
]
for i, (name, sub) in enumerate(enterprise):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(1.55 + row * 1.05)
    add_rect(s, x, y, Inches(3.95), Inches(0.92), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y + Inches(0.15), Inches(0.05), Inches(0.62), VIOLET)
    add_text(s, x + Inches(0.25), y + Inches(0.13), Inches(3.65), Inches(0.4),
             name, size=13, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.5), Inches(3.65), Inches(0.4),
             sub, size=9.5, color=SLATE)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 7 — Government & Public-Sector clients
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Government & Public Sector', 'Where our work touches the public infrastructure.', 7)

# Hero block — Sharjah
add_rect(s, Inches(0.6), Inches(1.55), Inches(6.1), Inches(2.7), NAVY)
add_text(s, Inches(0.85), Inches(1.7), Inches(5.7), Inches(0.4),
         'SHARJAH, UAE', size=10, bold=True, color=AMBER, tracking=3)
add_text(s, Inches(0.85), Inches(2.05), Inches(5.7), Inches(0.7),
         'Al Madam Municipality Portal',
         size=22, bold=True, color=WHITE)
add_text(s, Inches(0.85), Inches(2.7), Inches(5.7), Inches(1.4),
         'Citizen-services portal for Al Madam Municipality (Sharjah, UAE).\n'
         'React + TypeScript + Vite, deployed on Vercel.\n\n'
         'github.com/chatgptnotes/almadamonline · almadamonline.vercel.app',
         size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

# Hero block — Dubai
add_rect(s, Inches(6.85), Inches(1.55), Inches(6.1), Inches(2.7), INDIGO)
add_text(s, Inches(7.1), Inches(1.7), Inches(5.7), Inches(0.4),
         'DUBAI, UAE', size=10, bold=True, color=AMBER, tracking=3)
add_text(s, Inches(7.1), Inches(2.05), Inches(5.7), Inches(0.7),
         'Government of Dubai Media Office',
         size=22, bold=True, color=WHITE)
add_text(s, Inches(7.1), Inches(2.7), Inches(5.7), Inches(1.4),
         'JotForm Workflow Dashboard · Approval-Level Tracker.\n'
         'Sponsor: Huzaifa Dawasaz, Technical Project Manager (since 2022)\n'
         'and IT Infrastructure Engineer (since 2015) at the Office.\n\n'
         'Bettroi FZE delivery — operating from DTEC, Dubai Silicon Oasis.',
         size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

# Secondary government work
gov_more = [
    ('Manikaran Dam Safety',   'IoT compliance & automated reporting for Indian dam authorities'),
    ('IOC pilot',              'Indian Oil Corporation industrial-automation engagement'),
    ('NABH Online',            'NABH (national hospital accreditation) compliance toolchain'),
    ('Yellow Fever Ledger',    'Vaccine certification & ledger with government format compliance'),
    ('Sterling Tender Gen.',   'Government-tender document generator for Sterling Electricals'),
    ('ESIC AdamRIT',           'Patient-mapping system for ESIC scheme'),
]
add_text(s, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4),
         'Other public-sector work', size=12, bold=True, color=SLATE)
for i, (name, sub) in enumerate(gov_more):
    col = i % 3; row = i // 3
    x = Inches(0.6 + col * 4.15)
    y = Inches(4.95 + row * 1.0)
    add_rect(s, x, y, Inches(3.95), Inches(0.88), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y + Inches(0.15), Inches(0.05), Inches(0.6), AMBER)
    add_text(s, x + Inches(0.25), y + Inches(0.13), Inches(3.65), Inches(0.4),
             name, size=12, bold=True, color=NAVY)
    add_text(s, x + Inches(0.25), y + Inches(0.46), Inches(3.65), Inches(0.4),
             sub, size=9, color=SLATE)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 8 — Hospital & Enterprise clients
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Hospital & Enterprise Clients',
            'Self-operated hospitals + enterprises served via Bettroi & DrM Hope.', 8)

# Owned-and-operated
add_rect(s, Inches(0.6), Inches(1.55), Inches(6.1), Inches(3.4),
         WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
add_text(s, Inches(0.85), Inches(1.75), Inches(5.7), Inches(0.4),
         'OPERATED BY THE GROUP', size=10, bold=True, color=GREEN, tracking=3)
add_text(s, Inches(0.85), Inches(2.1), Inches(5.7), Inches(0.5),
         'Two NABH-accredited hospitals',
         size=18, bold=True, color=NAVY)
operated = [
    ('Hope Hospital, Nagpur',
     'Super-specialty · NABH-accredited · operating since 2005'),
    ('Ayushman Nagpur Hospital',
     'Orthopedic & spine surgery centre · founded 2018'),
]
for i, (name, sub) in enumerate(operated):
    y = Inches(2.85 + i * 0.95)
    add_text(s, Inches(0.85), y, Inches(5.7), Inches(0.4),
             name, size=14, bold=True, color=NAVY)
    add_text(s, Inches(0.85), y + Inches(0.35), Inches(5.7), Inches(0.55),
             sub, size=11, color=SLATE)

# Software clients (named or strongly inferred)
add_rect(s, Inches(6.85), Inches(1.55), Inches(6.1), Inches(3.4),
         WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
add_text(s, Inches(7.1), Inches(1.75), Inches(5.7), Inches(0.4),
         'NAMED SOFTWARE & ADVISORY CLIENTS', size=10, bold=True, color=BLUE, tracking=3)
named = [
    'ESIC (Employees\' State Insurance) — AdamRIT patient-mapping',
    'Sterling Electricals — auto tender doc generator',
    'Indian Oil Corporation — industrial-automation pilot',
    'Linkist (NFC smart cards) — AI back-office platform',
    'Justdial (Hope Hospital pipeline) — lead-management integration',
    'Vivah GMC — wedding planning portal (vivahgmc.com)',
    '4cbz, AIINMail, AutopanelDesign — Bettroi-built SaaS sites',
    'Bettroi disclosed: 50+ enterprise clients (GCC + India)',
]
add_bullets(s, Inches(7.1), Inches(2.15), Inches(5.7), Inches(2.7),
            named, size=11, gap=4)

# Disclosure callout
add_rect(s, Inches(0.6), Inches(5.2), Inches(12.4), Inches(1.6),
         RGBColor(0xFE, 0xF3, 0xC7), line=RGBColor(0xFD, 0xE6, 0x8A))
add_text(s, Inches(0.85), Inches(5.35), Inches(12), Inches(0.4),
         'A NOTE ON DISCLOSURE', size=10, bold=True, color=RGBColor(0x92, 0x40, 0x0E), tracking=3)
add_text(s, Inches(0.85), Inches(5.7), Inches(12), Inches(1),
         'Several enterprise engagements are under NDA and are not named here. '
         'The "50+ enterprise clients" figure on bettroi.com aggregates Bettroi\'s GCC + India '
         'roster across automation, healthcare, and oil-and-gas verticals. We can share specifics '
         'in a private 1:1 once a mutual NDA is in place.',
         size=11, color=RGBColor(0x78, 0x35, 0x0F))


# ═════════════════════════════════════════════════════════════════════════════
# Slide 9 — Engineering footprint (GitHub)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'Engineering Footprint',
            '133 public repositories on github.com/chatgptnotes — a small slice of what we build.', 9)

groups = [
    ('Hospital & Healthcare', GREEN, [
        'adamrit · adamrit-nextjs · adamrit_saas · adamrit-legacy',
        'drmhope.com-ESIC · drmhope-multitenancy · New_HMIS_Next_js',
        'aide · onepatientonedoctorwork · clinivoice · nabh.online',
        'ayushman-hospital (×4 sites) · hopehospital.com',
        'yellowfever (vaccine ledger) · one-scan-one-life',
        'RaftaarHealth · DigiHealthtwin · justdialwork',
    ]),
    ('Industrial / SCADA / Energy', AMBER, [
        'plcautopilot.com · plcautopilot-nextjs_with_AI',
        'EcoLogic-PLC-Assistant · autopaneldesign.com',
        'factorypulse.site · scada.work',
        'manikaran-dam-safety · sterling-tender-gen',
        'VisionClawRN · openclaw-for-hostinger',
    ]),
    ('Enterprise AI / Agents', VIOLET, [
        'agentsdr · zeroriskagent.com',
        'Linkist-01 · indiamonitorapp · moving-estimator',
        'aiinmail.com · economystic.com · 2menco · 4cbz.com',
        'jotform-dashboard · onepagebuild',
        'ironbark (open-source skill-harvest loop)',
    ]),
    ('Government / Public', BLUE, [
        'jotform-dashboard — Government of Dubai Media Office (UAE)',
        'almadamonline — Al Madam Municipality, Sharjah (UAE)',
        'manikaran-dam-safety — Indian dam compliance',
        'ioc — Indian Oil Corp engagement',
        'nabh — National accreditation toolchain',
    ]),
]
for i, (label, color, items) in enumerate(groups):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.55 + row * 2.7)
    add_rect(s, x, y, Inches(6), Inches(2.55), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(6), Inches(0.07), color)
    add_text(s, x + Inches(0.25), y + Inches(0.18), Inches(5.5), Inches(0.4),
             label, size=13, bold=True, color=NAVY)
    add_bullets(s, x + Inches(0.25), y + Inches(0.65), Inches(5.5), Inches(1.85),
                items, size=10, gap=2)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 10 — Tech stack
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_chrome(s, 'How We Build', 'Standard stack across the portfolio — fast to ship, easy to hand over.', 10)

stack = [
    ('AI & ML',           'GPT-4 / Claude / Gemini · LangChain · custom agents · RAG',  VIOLET),
    ('Web & App',         'Next.js · React · TypeScript · Tailwind · Vercel',           BLUE),
    ('Data',              'Supabase · PostgreSQL · pgvector · Prisma',                  INDIGO),
    ('Industrial',        'Modbus / OPC-UA · Node-RED · custom PLC bridges',           AMBER),
    ('Mobile',            'React Native (VisionClaw) · iOS / Android emergency apps',   PINK),
    ('DevOps',            'Hostinger VPS · Docker · GitHub Actions · CI/CD pipelines',  GREEN),
]
for i, (label, line, color) in enumerate(stack):
    col = i % 2; row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.7 + row * 1.65)
    add_rect(s, x, y, Inches(6), Inches(1.5), WHITE, line=RGBColor(0xE2, 0xE8, 0xF0))
    add_rect(s, x, y, Inches(0.18), Inches(1.5), color)
    add_text(s, x + Inches(0.45), y + Inches(0.25), Inches(5.4), Inches(0.5),
             label, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.45), y + Inches(0.78), Inches(5.4), Inches(0.65),
             line, size=12, color=SLATE)


# ═════════════════════════════════════════════════════════════════════════════
# Slide 11 — Contact
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
add_rect(s, 0, 0, Inches(4.5), Inches(0.18), BLUE)
add_rect(s, Inches(4.5), 0, Inches(4.5), Inches(0.18), VIOLET)
add_rect(s, Inches(9), 0, Inches(4.4), Inches(0.18), PINK)

add_text(s, Inches(0.8), Inches(1.4), Inches(11), Inches(0.4),
         'CONTACT · 121 BNI', size=11, bold=True, color=AMBER, tracking=3)
add_text(s, Inches(0.8), Inches(1.95), Inches(11), Inches(1.5),
         'Let\'s talk.',
         size=58, bold=True, color=WHITE)

# Cards
def contact_card(x, label, name, lines):
    add_rect(s, x, Inches(3.7), Inches(5.85), Inches(2.4),
             RGBColor(0x1E, 0x29, 0x3B), line=RGBColor(0x33, 0x41, 0x55))
    add_text(s, x + Inches(0.3), Inches(3.85), Inches(5.5), Inches(0.4),
             label, size=10, bold=True, color=AMBER, tracking=3)
    add_text(s, x + Inches(0.3), Inches(4.25), Inches(5.5), Inches(0.5),
             name, size=20, bold=True, color=WHITE)
    for i, line in enumerate(lines):
        add_text(s, x + Inches(0.3), Inches(4.85 + i * 0.32), Inches(5.5), Inches(0.32),
                 line, size=11, color=RGBColor(0xCB, 0xD5, 0xE1))

contact_card(Inches(0.8),  'INDIA OFFICE', 'Dr. BK Murali',
             ['Founder & CTO',
              '+91 93731 11709',
              'cmd@hopehospital.com',
              'hopetech.me · drmhope.com'])
contact_card(Inches(6.7),  'UAE OFFICE',   'Biji Thomas',
             ['CEO, Bettroi FZE',
              'A5, Techno-Hub, DTEC',
              'Dubai Silicon Oasis, UAE',
              'bettroi.com · hopephoenix.in'])

add_text(s, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
         'Card link · hopetech.me/bni/my-card.html', size=12, color=SLATE_LT)


# ═════════════════════════════════════════════════════════════════════════════
out = f'/Users/murali/1backup/BNI 121/HopeTech_Portfolio_{date.today().isoformat()}.pptx'
prs.save(out)
print('Saved:', out)
print('Slides:', len(prs.slides))
